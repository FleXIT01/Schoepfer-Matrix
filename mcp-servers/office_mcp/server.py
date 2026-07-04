"""
office_mcp — Word / PowerPoint / Excel (N6)
FastMCP-Server: create_docx, create_pptx, create_xlsx, read_docx, read_xlsx, read_pptx
Braucht: pip install python-docx python-pptx openpyxl
"""
from mcp.server.fastmcp import FastMCP
import os, json
from pathlib import Path
from datetime import datetime

mcp = FastMCP("office")

_OUTPUT_DIR = Path(os.environ.get("OFFICE_OUTPUT_DIR", r"n:\allinall\openclaw-workspace\output"))


def _ensure_output_dir():
    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def _default_out(stem: str, ext: str) -> str:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    return str(_OUTPUT_DIR / f"{stem}_{ts}.{ext}")


def _docx_ok():
    try:
        import docx
        return True
    except ImportError:
        return False


def _pptx_ok():
    try:
        import pptx
        return True
    except ImportError:
        return False


def _xl_ok():
    try:
        import openpyxl
        return True
    except ImportError:
        return False


def _md_to_docx(doc, content: str):
    """Konvertiert einfaches Markdown zu python-docx-Elementen."""
    from docx.shared import Pt
    for line in content.splitlines():
        stripped = line.strip()
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped[3:], level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped[2:], level=1)
        elif stripped.startswith(("- ", "* ", "• ")):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        elif stripped.startswith(("1. ", "2. ", "3. ", "4. ", "5. ")):
            doc.add_paragraph(stripped[3:], style="List Number")
        elif stripped == "---" or stripped == "***":
            doc.add_paragraph("─" * 40)
        elif stripped == "":
            doc.add_paragraph("")
        else:
            doc.add_paragraph(stripped)


@mcp.tool()
def create_docx(title: str, content: str, output_file: str = "") -> str:
    """Erstellt ein Word-Dokument (.docx) aus Titel und Markdown-Inhalt.

    Args:
        title: Dokumenttitel (wird als H1-Überschrift eingefügt)
        content: Inhalt als Markdown (# H1, ## H2, - Liste, normaler Text)
        output_file: Zieldatei-Pfad (leer = automatisch in output/)

    Returns:
        Pfad zur erstellten .docx-Datei.
    """
    if not _docx_ok():
        return "[office] python-docx nicht installiert. pip install python-docx"

    from docx import Document

    _ensure_output_dir()
    out = output_file.strip() or _default_out("dokument", "docx")

    doc = Document()
    doc.add_heading(title, level=0)
    _md_to_docx(doc, content)

    Path(out).parent.mkdir(parents=True, exist_ok=True)
    doc.save(out)
    return f"[create_docx] OK {out}"


@mcp.tool()
def read_docx(file_path: str) -> str:
    """Liest ein Word-Dokument (.docx) und gibt den Text als Markdown zurück.

    Args:
        file_path: Pfad zur .docx-Datei

    Returns:
        Dokumentinhalt als Text.
    """
    if not _docx_ok():
        return "[office] python-docx nicht installiert. pip install python-docx"

    from docx import Document

    p = Path(file_path)
    if not p.exists():
        return f"[read_docx] Datei nicht gefunden: {file_path}"

    doc = Document(str(p))
    lines = []
    for para in doc.paragraphs:
        style = para.style.name.lower()
        text = para.text.strip()
        if not text:
            lines.append("")
            continue
        if "heading 1" in style:
            lines.append(f"# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
        elif "heading 3" in style:
            lines.append(f"### {text}")
        elif "list bullet" in style:
            lines.append(f"- {text}")
        elif "list number" in style:
            lines.append(f"1. {text}")
        else:
            lines.append(text)
    return "\n".join(lines)


@mcp.tool()
def create_pptx(title: str, slides_json: str, output_file: str = "") -> str:
    """Erstellt eine PowerPoint-Präsentation (.pptx).

    Args:
        title: Präsentationstitel (Titelfolie)
        slides_json: JSON-Array mit Folien: [{"title":"...", "content":"Bullet1\\nBullet2"}]
        output_file: Zieldatei-Pfad (leer = automatisch in output/)

    Returns:
        Pfad zur erstellten .pptx-Datei.
    """
    if not _pptx_ok():
        return "[office] python-pptx nicht installiert. pip install python-pptx"

    from pptx import Presentation
    from pptx.util import Inches, Pt

    try:
        slides_data = json.loads(slides_json)
    except json.JSONDecodeError as e:
        return f"[create_pptx] Ungültiges JSON in slides_json: {e}"

    if not isinstance(slides_data, list):
        return "[create_pptx] slides_json muss ein JSON-Array sein."

    _ensure_output_dir()
    out = output_file.strip() or _default_out("praesentation", "pptx")

    prs = Presentation()

    # Titelfolie
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = datetime.now().strftime("%d.%m.%Y")

    # Inhaltsfolien
    for s in slides_data:
        slide_title = s.get("title", "")
        slide_content = s.get("content", "")
        layout = prs.slide_layouts[1]  # Title and Content
        sl = prs.slides.add_slide(layout)
        sl.shapes.title.text = slide_title
        tf = sl.placeholders[1].text_frame
        tf.clear()
        lines = slide_content.split("\n")
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            if line.startswith(("- ", "* ", "• ")):
                line = line[2:]
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
                p.level = 1

    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return f"[create_pptx] OK {out}"


@mcp.tool()
def read_pptx(file_path: str) -> str:
    """Liest eine PowerPoint-Datei und gibt Text je Folie zurück.

    Args:
        file_path: Pfad zur .pptx-Datei

    Returns:
        Text aller Folien, durchnummeriert.
    """
    if not _pptx_ok():
        return "[office] python-pptx nicht installiert. pip install python-pptx"

    from pptx import Presentation

    p = Path(file_path)
    if not p.exists():
        return f"[read_pptx] Datei nicht gefunden: {file_path}"

    prs = Presentation(str(p))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        lines.append(f"## Folie {i}\n" + "\n".join(texts))
    return "\n\n".join(lines)


@mcp.tool()
def create_xlsx(
    data_json: str,
    output_file: str = "",
    sheet_name: str = "Tabelle1",
) -> str:
    """Erstellt eine Excel-Datei (.xlsx) aus JSON-Daten.

    Args:
        data_json: JSON-Array — entweder Liste von Dicts (Schlüssel=Spalten) oder
                   Liste von Listen (erste Zeile = Header). Beispiel:
                   '[{"Name":"Max","Wert":42},{"Name":"Eva","Wert":99}]'
        output_file: Zieldatei-Pfad (leer = automatisch in output/)
        sheet_name: Name des Arbeitsblatts (Default: Tabelle1)

    Returns:
        Pfad zur erstellten .xlsx-Datei.
    """
    if not _xl_ok():
        return "[office] openpyxl nicht installiert. pip install openpyxl"

    import openpyxl
    from openpyxl.styles import Font

    try:
        data = json.loads(data_json)
    except json.JSONDecodeError as e:
        return f"[create_xlsx] Ungültiges JSON: {e}"

    if not isinstance(data, list) or not data:
        return "[create_xlsx] data_json muss ein nicht-leeres JSON-Array sein."

    _ensure_output_dir()
    out = output_file.strip() or _default_out("tabelle", "xlsx")

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name[:31]

    if isinstance(data[0], dict):
        headers = list(data[0].keys())
        ws.append(headers)
        for cell in ws[1]:
            cell.font = Font(bold=True)
        for row_dict in data:
            ws.append([row_dict.get(h, "") for h in headers])
    elif isinstance(data[0], list):
        for row in data:
            ws.append(row)
        for cell in ws[1]:
            cell.font = Font(bold=True)
    else:
        return "[create_xlsx] data_json-Elemente müssen Dicts oder Listen sein."

    Path(out).parent.mkdir(parents=True, exist_ok=True)
    wb.save(out)
    return f"[create_xlsx] OK {out}"


@mcp.tool()
def read_xlsx(file_path: str, sheet_name: str = "") -> str:
    """Liest eine Excel-Datei und gibt den Inhalt als CSV-ähnlichen Text zurück.

    Args:
        file_path: Pfad zur .xlsx-Datei
        sheet_name: Name des Arbeitsblatts (leer = erstes Blatt)

    Returns:
        Tabelleninhalt als Tabulator-getrennte Zeilen (erste Zeile = Header).
    """
    if not _xl_ok():
        return "[office] openpyxl nicht installiert. pip install openpyxl"

    import openpyxl

    p = Path(file_path)
    if not p.exists():
        return f"[read_xlsx] Datei nicht gefunden: {file_path}"

    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active

    rows = []
    for row in ws.iter_rows(values_only=True):
        rows.append("\t".join(str(c) if c is not None else "" for c in row))
    wb.close()

    if not rows:
        return "[read_xlsx] Leeres Arbeitsblatt."
    return f"[read_xlsx] {ws.title} ({len(rows)} Zeilen)\n" + "\n".join(rows[:200])


if __name__ == "__main__":
    mcp.run(transport="stdio")
